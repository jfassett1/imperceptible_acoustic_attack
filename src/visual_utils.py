from whisper import log_mel_spectrogram
import matplotlib.pyplot as plt
import numpy as np
import torch
import librosa
from pathlib import Path
from scipy.io import wavfile
import whisper

sample_raw = whisper.load_audio("/home/jaydenfassett/audioversarial/imperceptible/original_audio.wav")


def show(noise,
         sample_audio,
         PATHS,
         prepend:bool=False):



    #Saving Audio
    # audio_list = [sound for sound in (PATHS.example_dir / "sample_sounds").glob("*.wav")] # List of paths
    # images_list,audio_list = audio_to_img(noise,audio_list,sample_audio,PATHS.audio_dir)

    if prepend:
        save_photo_prepend(noise,sample_audio,PATHS.img_dir/"plot.png")
    else:
        save_photo_overlay(noise,sample_audio,PATHS.img_dir/"plot.png")
    print(f"Saving image to {PATHS.img_dir/'plot.png'}")

def raw_to_mel(x,device="cpu"):
    x = whisper.pad_or_trim(x)
    mel = whisper.log_mel_spectrogram(x).to(device).unsqueeze(dim=0)
    return mel


def overlay(noise, raw_audio):
    if noise.ndim == 1:
        noise = noise[np.newaxis, :]
    if raw_audio.ndim == 1:
        raw_audio = raw_audio[np.newaxis, :]

    if noise.shape[-1] > raw_audio.shape[-1]:
        noise = noise[:, :raw_audio.shape[-1]]

    pad_length = raw_audio.shape[-1] - noise.shape[-1]
    padding = np.zeros((noise.shape[0], pad_length), dtype=noise.dtype)
    noise_padded = np.concatenate([noise, padding], axis=-1)
    result = noise_padded + raw_audio

    if result.shape[0] == 1:
        return result.squeeze(0)
    return result

def save_photo_overlay(noise, audio, savepath):
    """
    Save an overlay of the noise and audio spectrogram as an image.

    Parameters:
    noise (np.ndarray): The noise array.
    audio (np.ndarray): The audio array.
    savepath (str): Path where the image will be saved.
    """
    # Ensure noise and audio have the same length

    noise = np.concatenate([noise.squeeze(), np.zeros(len(audio) - len(noise))])

    # Combine noise and audio
    combined_audio = noise + audio

    hop_length = 160
    mel_spec1 = log_mel_spectrogram(audio.astype(np.float32)).numpy()
    mel_spec2 = log_mel_spectrogram(combined_audio.astype(np.float32)).numpy()

    # mel_spectrogram1 = librosa.feature.melspectrogram(y=aud1, sr=fs1, n_mels=80, hop_length=hop_length)
    # mel_spectrogram_db1 = librosa.power_to_db(mel_spectrogram1, ref=np.max)

    # mel_spectrogram2 = librosa.feature.melspectrogram(y=aud2, sr=fs1, n_mels=80, hop_length=hop_length)
    # mel_spectrogram_db2 = librosa.power_to_db(mel_spectrogram2, ref=np.max)

    fig,axes = plt.subplots(1,2,figsize = (12,5))

    librosa.display.specshow(mel_spec1, sr=16000, hop_length=hop_length, 
                            x_axis='time', y_axis='mel', ax=axes[0],shading="gouraud",cmap="magma")
    
    librosa.display.specshow(mel_spec2, sr=16000, hop_length=hop_length, 
                            x_axis='time', y_axis='mel', ax=axes[1],shading="gouraud",cmap="magma")
    axes[0].set_title(f"Unperturbed Audio ({len(combined_audio)/16000:.2f} seconds)")
    axes[1].set_title(f"Perturbed Audio ({len(combined_audio)/16000:.2f} seconds)")
    # Save the figure
    plt.tight_layout()
    plt.savefig(savepath)
    plt.close(fig)


def save_photo_prepend(noise, audio, savepath):
    """
    Save prepended noise and audio spectrogram as an image.

    Parameters:
    noise (np.ndarray): The noise array.
    audio (np.ndarray): The audio array.
    savepath (str): Path where the image will be saved.
    """
    # Ensure noise and audio have the same length
    audio = audio.astype(np.float32)

    # noise = np.concatenate([noise.squeeze(), np.zeros(len(audio) - len(noise))])
    combined_audio = np.concatenate([noise.squeeze(),audio])
    # Combine noise and audio
    # combined_audio = noise + audio

    hop_length = 160
    mel_spec1 = log_mel_spectrogram(audio).numpy()
    mel_spec2 = log_mel_spectrogram(combined_audio.astype(np.float32)).numpy()

    # mel_spectrogram1 = librosa.feature.melspectrogram(y=aud1, sr=fs1, n_mels=80, hop_length=hop_length)
    # mel_spectrogram_db1 = librosa.power_to_db(mel_spectrogram1, ref=np.max)

    # mel_spectrogram2 = librosa.feature.melspectrogram(y=aud2, sr=fs1, n_mels=80, hop_length=hop_length)
    # mel_spectrogram_db2 = librosa.power_to_db(mel_spectrogram2, ref=np.max)

    fig,axes = plt.subplots(1,2,figsize = (12,5))

    librosa.display.specshow(mel_spec1, sr=16000, hop_length=hop_length, 
                            x_axis='time', y_axis='mel', ax=axes[0],shading="gouraud",cmap="magma")
    
    librosa.display.specshow(mel_spec2, sr=16000, hop_length=hop_length, 
                            x_axis='time', y_axis='mel', ax=axes[1],shading="gouraud",cmap="magma")
    axes[0].set_title(f"Unperturbed Audio ({len(combined_audio)/160000:.2f} seconds)")
    axes[1].set_title(f"Perturbed Audio ({len(combined_audio)/160000:.2f} seconds)")
    # Save the figure
    plt.tight_layout()
    plt.savefig(savepath)
    plt.close(fig)
sample_mel = raw_to_mel(sample_raw)

def contrast(aud1,aud2):

    fs1, aud1 = aud1
    fs2, aud2 = aud2
    aud1 = aud1.astype(np.float32)
    aud2 = aud2.astype(np.float32)

    hop_length = 160
    mel_spec1 = log_mel_spectrogram(aud1).numpy()
    mel_spec2 = log_mel_spectrogram(aud2).numpy()
    # mel_spectrogram1 = librosa.feature.melspectrogram(y=aud1, sr=fs1, n_mels=80, hop_length=hop_length)
    # mel_spectrogram_db1 = librosa.power_to_db(mel_spectrogram1, ref=np.max)

    # mel_spectrogram2 = librosa.feature.melspectrogram(y=aud2, sr=fs1, n_mels=80, hop_length=hop_length)
    # mel_spectrogram_db2 = librosa.power_to_db(mel_spectrogram2, ref=np.max)

    fig,axes = plt.subplots(1,2,figsize = (12,5))

    librosa.display.specshow(mel_spec1, sr=fs1, hop_length=hop_length, 
                            x_axis='time', y_axis='mel', ax=axes[0],shading="gouraud",cmap="magma")
    
    librosa.display.specshow(mel_spec2, sr=fs2, hop_length=hop_length, 
                            x_axis='time', y_axis='mel', ax=axes[1],shading="gouraud",cmap="magma")
    axes[0].set_title(f"Unperturbed Audio ({len(aud1)/fs1:.2f} seconds)")
    axes[1].set_title(f"Perturbed Audio ({len(aud2)/fs2:.2f} seconds)")
    return fig

def show_mel_spectrograms(mel_specs, sr=16000, hop_length=160, titles=None):
    num_specs = len(mel_specs)
    fig, axes = plt.subplots(1, num_specs, figsize=(5*num_specs, 5))
    
    if num_specs == 1:
        axes = [axes]
    
    for i, mel in enumerate(mel_specs):
        # Convert tensor to numpy array if necessary.
        if not isinstance(mel, np.ndarray):
            # For PyTorch tensors: detach and move to CPU if needed.
            if hasattr(mel, "detach"):
                mel = mel.detach().cpu().numpy()
            # For other tensors that already have a numpy() method.
            elif hasattr(mel, "numpy"):
                mel = mel.numpy()
            else:
                raise ValueError("Input mel_spec is not a numpy array or a convertible tensor.")
                
        ax = axes[i]
        librosa.display.specshow(mel, sr=sr, hop_length=hop_length,
                                 x_axis='time', y_axis='mel', ax=ax,
                                 shading="gouraud", cmap="magma")
        title = titles[i] if titles is not None and i < len(titles) else f"Spectrogram {i+1}"
        ax.set_title(title)
    
    plt.tight_layout()
    return fig

def audio_to_img(noise, #Learned noise
                 audio_list, #Sample noises to overlay over audio
                 raw_sample, #Raw audio sample to overlay audio over
                 save_dir: Path, #Directory for saving resultant images
                 sampling_rate=16000):
    """
    Takes in audio
    Returns list of paths to images
    """
    img_paths = []
    audio_paths = [] 
    audio_list.append(noise)

    for i,audio in enumerate(audio_list):
        if isinstance(audio, torch.Tensor):
            aud_path = "noise"
            audio = audio.detach().cpu().numpy()
        else:
            aud_path = audio
            audio,sr = librosa.load(audio,sr=None)
            audio = librosa.resample(audio, orig_sr=sr, target_sr=16000)
        audio = audio.astype(np.float32)
        # print(audio.mean(),raw_sample.mean(),aud_path)
        audio = normalize_volume(audio,raw_sample)
        audio = overlay(audio,raw_sample)
        hop_length = 160
        mel_spec1 = log_mel_spectrogram(audio).numpy().squeeze()

        plt.figure(figsize=(16, 5))  
        librosa.display.specshow(mel_spec1, sr=sampling_rate, hop_length=hop_length, 
                                x_axis='time', y_axis='mel', shading="gouraud", cmap="magma")
        plt.title(f"Unperturbed Audio ({len(audio)/sampling_rate:.2f} seconds)")
        plt.colorbar()
        image_path = save_dir / f"{i}.png"
        plt.savefig(image_path, bbox_inches='tight', pad_inches=0)
        img_paths.append(str(save_dir / f"{i}.png"))
        new_audio = librosa.resample(audio,orig_sr=16000,target_sr=44100)
        new_audio = (new_audio * 32767).astype(np.int16)
        # new_audio = np.int16((new_audio - 0.5) * 2 * 32767)
        np.save(save_dir / f"{i}.np",new_audio)

        wavfile.write(save_dir / f"{i}.wav",44100,new_audio)
        audio_paths.append(str(save_dir / f"{i}.np.npy"))
    # plt.colorbar(format="%+2.0f dB")
    return (img_paths,audio_paths)

from pptx import Presentation
from pptx.util import Inches, Pt


def normalize_volume(signal1, signal2):
    """
    Normalize the volume of signal1 to match the volume of signal2.

    Parameters:
    - signal1 (numpy.ndarray): The audio signal to be normalized.
    - signal2 (numpy.ndarray): The reference audio signal.

    Returns:
    - numpy.ndarray: The volume-normalized version of signal1.
    """
    # Ensure both signals are NumPy arrays
    signal1 = np.asarray(signal1)
    signal2 = np.asarray(signal2)

    # Calculate RMS (Root Mean Square) of both signals
    rms1 = np.sqrt(np.mean(signal1 ** 2))
    rms2 = np.sqrt(np.mean(signal2 ** 2))

    # Avoid division by zero
    if rms1 == 0:
        raise ValueError("The first signal (signal1) has zero volume (RMS). Cannot normalize.")

    # Compute the scaling factor
    scaling_factor = rms2 / rms1

    # Normalize signal1 to match the volume of signal2
    normalized_signal = signal1 * scaling_factor

    return normalized_signal

def generate_example_ppt(image_paths, audio_paths, output_path="stacked_images.pptx"):
    if not isinstance(image_paths[0],str):
        image_paths = [str(img) for img in image_paths]

    # Create a PowerPoint presentation
    prs = Presentation()

    # Set slide dimensions (optional, default is widescreen 16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    left_space = Inches(6)  # Leave space on the left for audio
    top_margin = Inches(0.5)  # Top margin
    spacing = Inches(0.5)  # Space between images
    img_height = (prs.slide_height - top_margin * 2 - spacing * 4) / 5  # Dynamic height for 5 images

    # Add images stacked vertically
    for i, (img_path,aud_path) in enumerate(zip(image_paths,audio_paths)):

        top = top_margin + i * (img_height + spacing)
        slide.shapes.add_picture(img_path, left_space, top, height=img_height) # Add images uniformly spaced vertically

        slide.shapes.add_movie(                                                # Add audio at same positions
            # r"/home/jaydenfassett/audioversarial/imperceptible/audio_with_attack.wav",
            aud_path,
            left=Inches(4.17),
            top=top,       
            width=Inches(1.67),
            height=Inches(0.76),
            poster_frame_image=None,
            mime_type='audio/mp3'
        )
                # Save the presentation
    prs.save(output_path)
    print(f"Presentation saved as {output_path}")


if __name__ == "__main__":
    print("")